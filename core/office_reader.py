"""
Office File Reader - Enhanced support for Microsoft Office files
"""

import os

class OfficeFileReader:
    """Enhanced office file reader with better error handling"""
    
    def __init__(self):
        self.supported_formats = {
            'docx': self.read_docx,
            'xlsx': self.read_xlsx,
            'pptx': self.read_pptx,
            'pdf': self.read_pdf
        }
    
    def read_docx(self, file_path):
        """Read Word document content"""
        try:
            from docx import Document
            doc = Document(file_path)
            content = ["📄 Word Document\n"]
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            # Read tables
            for table in doc.tables:
                content.append("\n📊 Table:")
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    content.append(" | ".join(row_data))
            
            return "\n".join(content)
            
        except ImportError:
            return "📄 Word Document\n\n⚠️ python-docx library not installed.\nInstall with: pip install python-docx"
        except Exception as e:
            return f"📄 Word Document\n\n❌ Error reading document: {str(e)}"
    
    def read_xlsx(self, file_path):
        """Read Excel spreadsheet content"""
        try:
            from openpyxl import load_workbook
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            content = ["📊 Excel Spreadsheet\n"]
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"\n📋 Sheet: {sheet_name}")
                content.append("─" * 30)
                
                for row_num, row in enumerate(sheet.iter_rows(max_row=20, values_only=True)):
                    if row_num > 20:
                        break
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):
                        content.append(" | ".join(row_data[:10]))  # First 10 columns
            
            workbook.close()
            return "\n".join(content)
            
        except ImportError:
            return "📊 Excel Spreadsheet\n\n⚠️ openpyxl library not installed.\nInstall with: pip install openpyxl"
        except Exception as e:
            return f"📊 Excel Spreadsheet\n\n❌ Error reading spreadsheet: {str(e)}"
    
    def read_pptx(self, file_path):
        """Read PowerPoint presentation content"""
        try:
            from pptx import Presentation
            presentation = Presentation(file_path)
            content = ["📽️ PowerPoint Presentation\n"]
            content.append(f"📊 Total Slides: {len(presentation.slides)}")
            content.append("─" * 30)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                content.append(f"\n🔸 Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(f"  • {shape.text.strip()}")
            
            return "\n".join(content)
            
        except ImportError:
            return "📽️ PowerPoint Presentation\n\n⚠️ python-pptx library not installed.\nInstall with: pip install python-pptx"
        except Exception as e:
            return f"📽️ PowerPoint Presentation\n\n❌ Error reading presentation: {str(e)}"
    
    def read_pdf(self, file_path):
        """Read PDF document content"""
        try:
            import PyPDF2
            content = []
            
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                content.append("📕 PDF Document")
                content.append(f"📊 Total Pages: {len(reader.pages)}")
                content.append("─" * 30)
                
                # Extract text from first few pages
                for page_num, page in enumerate(reader.pages[:5], 1):
                    try:
                        page_text = page.extract_text()
                        
                        if page_text.strip():
                            content.append(f"\n📄 Page {page_num}:")
                            content.append(page_text.strip()[:1000])  # First 1000 chars
                        else:
                            content.append(f"\n📄 Page {page_num}: (No extractable text)")
                    except Exception as e:
                        content.append(f"\n📄 Page {page_num}: Error extracting text - {str(e)}")
                
                if len(reader.pages) > 5:
                    content.append(f"\n... and {len(reader.pages) - 5} more pages")
            
            return "\n".join(content)
            
        except ImportError:
            return "📕 PDF Document\n\n⚠️ PyPDF2 library not installed.\nInstall with: pip install PyPDF2"
        except Exception as e:
            return f"📕 PDF Document\n\n❌ Error reading PDF: {str(e)}"

    def can_read_format(self, file_extension):
        """Check if the reader can handle this format"""
        return file_extension.lower().lstrip('.') in self.supported_formats

    def read_file(self, file_path):
        """Generic file reader that determines format automatically"""
        file_ext = os.path.splitext(file_path)[1].lower().lstrip('.')
        
        if file_ext in self.supported_formats:
            return self.supported_formats[file_ext](file_path)
        else:
            return f"📄 {file_ext.upper()} File\n\n⚠️ Unsupported file format for office preview"